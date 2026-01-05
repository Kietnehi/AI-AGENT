"""Local LLM Tool with Qwen 2.5B and Gemini API Support"""
from transformers import AutoTokenizer, AutoModelForCausalLM, pipeline, BitsAndBytesConfig
import torch
from typing import Optional, Dict, Any
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from diffusers import StableDiffusionPipeline
from PIL import Image
import os
import json
import random

class LocalLLM:
    """Local Language Model using Qwen 2.5B with 4-bit quantization"""
    
    def __init__(self, model_name: str = "Qwen/Qwen2.5-1.5B-Instruct"):
        """
        Initialize Local LLM with 4-bit quantization for low VRAM (4GB GPU)
        
        Args:
            model_name: Hugging Face model name (default: Qwen2.5-1.5B-Instruct for lighter weight)
        """
        self.model_name = model_name
        self.model = None
        self.tokenizer = None
        self.pipeline = None
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"🖥️ Local LLM will use device: {self.device}")
        if self.device == "cuda":
            print(f"💾 GPU VRAM optimization: 4-bit quantization enabled")
    
    def load_model(self):
        """Lazy load the model with 4-bit quantization"""
        if self.model is None:
            try:
                print(f"🔄 Loading local LLM: {self.model_name}...")
                
                # Load tokenizer
                self.tokenizer = AutoTokenizer.from_pretrained(
                    self.model_name,
                    trust_remote_code=True
                )
                
                # Configure 4-bit quantization for 4GB VRAM
                if self.device == "cuda":
                    quantization_config = BitsAndBytesConfig(
                        load_in_4bit=True,
                        bnb_4bit_compute_dtype=torch.float16,
                        bnb_4bit_use_double_quant=True,
                        bnb_4bit_quant_type="nf4"
                    )
                    
                    print("⚙️ Using 4-bit quantization (NF4) for low VRAM...")
                    
                    # Load model with 4-bit quantization
                    self.model = AutoModelForCausalLM.from_pretrained(
                        self.model_name,
                        quantization_config=quantization_config,
                        device_map="auto",
                        trust_remote_code=True
                    )
                else:
                    # CPU mode - no quantization
                    self.model = AutoModelForCausalLM.from_pretrained(
                        self.model_name,
                        torch_dtype=torch.float32,
                        trust_remote_code=True
                    )
                    self.model = self.model.to(self.device)
                
                # Create pipeline without device parameter when using device_map="auto"
                self.pipeline = pipeline(
                    "text-generation",
                    model=self.model,
                    tokenizer=self.tokenizer
                )
                
                print("✓ Local LLM loaded successfully with 4-bit quantization!")
                print("📊 Memory usage: ~1-1.5GB VRAM (model weights)")
                
            except Exception as e:
                print(f"✗ Error loading local LLM: {e}")
                raise
        
        return self.pipeline
    
    def generate(
        self,
        prompt: str,
        max_length: int = 512,
        temperature: float = 0.7,
        top_p: float = 0.9,
        top_k: int = 50
    ) -> Dict[str, Any]:
        """
        Generate text using local LLM
        
        Args:
            prompt: Input text prompt
            max_length: Maximum length of generated text
            temperature: Sampling temperature (higher = more random)
            top_p: Nucleus sampling parameter
            top_k: Top-k sampling parameter
            
        Returns:
            Dict with generated text and metadata
        """
        try:
            # Load model if not loaded
            pipe = self.load_model()
            
            # Format prompt for Qwen chat model
            messages = [
                {"role": "user", "content": prompt}
            ]
            
            # Generate
            outputs = pipe(
                messages,
                max_new_tokens=max_length,
                temperature=temperature,
                top_p=top_p,
                top_k=top_k,
                do_sample=True,
                return_full_text=False
            )
            
            # Extract generated text
            generated_text = outputs[0]["generated_text"]
            
            return {
                "success": True,
                "response": generated_text,
                "model": self.model_name,
                "device": self.device
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "model": self.model_name
            }
    
    def chat(
        self,
        message: str,
        history: Optional[list] = None,
        max_length: int = 512,
        temperature: float = 0.7
    ) -> Dict[str, Any]:
        """
        Chat with the model (with conversation history)
        
        Args:
            message: User message
            history: Conversation history as list of {"role": "user"/"assistant", "content": "..."}
            max_length: Maximum length of response
            temperature: Sampling temperature
            
        Returns:
            Dict with response and updated history
        """
        try:
            # Load model if not loaded
            pipe = self.load_model()
            
            # Build messages
            messages = history if history else []
            messages.append({"role": "user", "content": message})
            
            # Generate
            outputs = pipe(
                messages,
                max_new_tokens=max_length,
                temperature=temperature,
                do_sample=True,
                return_full_text=False
            )
            
            # Extract response
            response = outputs[0]["generated_text"]
            
            # Update history
            messages.append({"role": "assistant", "content": response})
            
            return {
                "success": True,
                "response": response,
                "history": messages,
                "model": self.model_name
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "model": self.model_name
            }
    
    def create_presentation_slides(self, topic: str, num_slides: int = 5) -> Dict[str, Any]:
        """
        Create presentation slides using Qwen local LLM with AI-generated images
        
        Args:
            topic: Topic for the presentation
            num_slides: Number of slides to create
            
        Returns:
            Dict with slide content and file path
        """
        try:
            # Generate slide content using Qwen
            prompt = f"""Tạo một bài thuyết trình về chủ đề: {topic}

Yêu cầu:
- Tạo {num_slides} slide
- Mỗi slide có tiêu đề và nội dung chi tiết
- Nội dung phải logic, mạch lạc
- Sử dụng bullet points
- Thêm trường "image_prompt" cho một số slide để tạo hình minh họa (bằng tiếng Anh, mô tả chi tiết)
- Format JSON như sau:
{{
  "title": "Tiêu đề bài thuyết trình",
  "slides": [
    {{
      "title": "Tiêu đề slide",
      "content": ["Điểm 1", "Điểm 2", "Điểm 3"],
      "image_prompt": "detailed image description in English (optional)"
    }}
  ]
}}

Chỉ trả về JSON, không thêm text khác."""

            # Load model if not loaded
            pipe = self.load_model()
            
            # Generate
            messages = [{"role": "user", "content": prompt}]
            outputs = pipe(
                messages,
                max_new_tokens=1024,
                temperature=0.7,
                do_sample=True,
                return_full_text=False
            )
            
            # Parse JSON response
            text_response = outputs[0]["generated_text"]
            # Clean up markdown code blocks if present
            if "```json" in text_response:
                text_response = text_response.split("```json")[1].split("```")[0].strip()
            elif "```" in text_response:
                text_response = text_response.split("```")[1].split("```")[0].strip()
            
            slides_data = json.loads(text_response)
            
            # Initialize Text-to-Image model
            text_to_image = TextToImage()
            
            # Determine which slides should have images (about 60% of content slides)
            content_slides = slides_data.get("slides", [])
            num_images = min(int(num_slides * 0.6), len(content_slides))
            
            # Create PowerPoint presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = slides_data.get("title", topic)
            
            # Content slides
            images_generated = 0
            for idx, slide_data in enumerate(content_slides):
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                title_shape = slide.shapes.title
                title_shape.text = slide_data.get("title", "")
                
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                
                for point in slide_data.get("content", []):
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                
                # Generate and add image if needed
                should_add_image = (
                    images_generated < num_images and
                    (slide_data.get("image_prompt") or idx % 2 == 0)  # Add image to even slides or if prompt exists
                )
                
                if should_add_image:
                    # Use provided prompt or generate one from slide title
                    image_prompt = slide_data.get("image_prompt") or f"{slide_data.get('title', topic)}, professional illustration, high quality"
                    
                    # Generate image
                    result = text_to_image.generate_image(image_prompt)
                    
                    if result["success"]:
                        # Add image to slide (bottom right corner)
                        left = Inches(7)  # Right side
                        top = Inches(4.5)  # Bottom
                        width = Inches(2.5)
                        
                        slide.shapes.add_picture(
                            result["image_path"],
                            left, top,
                            width=width
                        )
                        images_generated += 1
                        print(f"✓ Added image to slide {idx + 1}")
            
            # Cleanup Text-to-Image model
            text_to_image.cleanup()
            
            # Save presentation
            os.makedirs("slides", exist_ok=True)
            filename = f"slides/presentation_{topic[:30].replace(' ', '_')}.pptx"
            prs.save(filename)
            
            return {
                "success": True,
                "filename": filename,
                "title": slides_data.get("title", topic),
                "num_slides": len(content_slides) + 1,
                "num_images": images_generated,
                "model": self.model_name
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "model": self.model_name
            }
    
    def cleanup(self):
        """Clean up model from memory"""
        if self.model is not None:
            del self.model
            self.model = None
        
        if self.tokenizer is not None:
            del self.tokenizer
            self.tokenizer = None
        
        if self.pipeline is not None:
            del self.pipeline
            self.pipeline = None
        
        # Clear CUDA cache
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
        
        print("Local LLM cleaned up")


class TextToImage:
    """Text-to-Image generation using Stable Diffusion"""
    
    def __init__(self, model_id: str = "prompthero/openjourney"):
        """
        Initialize Text-to-Image model
        
        Args:
            model_id: Hugging Face model ID (default: prompthero/openjourney)
        """
        self.model_id = model_id
        self.pipe = None
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"🖼️ Text-to-Image will use device: {self.device}")
    
    def load_model(self):
        """Lazy load the Stable Diffusion model"""
        if self.pipe is None:
            try:
                print(f"🔄 Loading Text-to-Image model: {self.model_id}...")
                
                if self.device == "cuda":
                    self.pipe = StableDiffusionPipeline.from_pretrained(
                        self.model_id,
                        torch_dtype=torch.float16
                    )
                    self.pipe = self.pipe.to("cuda")
                    print("✓ Text-to-Image model loaded successfully on GPU!")
                else:
                    self.pipe = StableDiffusionPipeline.from_pretrained(
                        self.model_id,
                        torch_dtype=torch.float32
                    )
                    self.pipe = self.pipe.to("cpu")
                    print("✓ Text-to-Image model loaded successfully on CPU!")
                    print("⚠️ Warning: CPU inference will be slow. Consider using GPU for faster generation.")
                
            except Exception as e:
                print(f"✗ Error loading Text-to-Image model: {e}")
                raise
        
        return self.pipe
    
    def generate_image(self, prompt: str, output_path: str = None) -> Dict[str, Any]:
        """
        Generate image from text prompt
        
        Args:
            prompt: Text description for image generation
            output_path: Path to save the image (optional)
            
        Returns:
            Dict with image path and metadata
        """
        try:
            # Load model if not loaded
            pipe = self.load_model()
            
            print(f"🎨 Generating image: '{prompt[:50]}...'")
            
            # Generate image
            image = pipe(prompt).images[0]
            
            # Save image
            if output_path is None:
                os.makedirs("slides/images", exist_ok=True)
                output_path = f"slides/images/generated_{random.randint(1000, 9999)}.png"
            
            image.save(output_path)
            print(f"✓ Image saved: {output_path}")
            
            return {
                "success": True,
                "image_path": output_path,
                "prompt": prompt,
                "model": self.model_id
            }
            
        except Exception as e:
            print(f"✗ Error generating image: {e}")
            return {
                "success": False,
                "error": str(e),
                "model": self.model_id
            }
    
    def cleanup(self):
        """Clean up model from memory"""
        if self.pipe is not None:
            del self.pipe
            self.pipe = None
        
        # Clear CUDA cache
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
        
        print("Text-to-Image model cleaned up")


class GeminiAPI:
    """Gemini API LLM integration"""
    
    def __init__(self, api_key: str, model_name: str = "gemini-2.5-flash"):
        """
        Initialize Gemini API
        
        Args:
            api_key: Gemini API key
            model_name: Model name (default: gemini-2.5-flash)
        """
        self.api_key = api_key
        self.model_name = model_name
        self.client = genai.Client(api_key=api_key)
        print(f"✓ Gemini API initialized with model: {model_name}")
    
    def generate(self, prompt: str, max_length: int = 2048, temperature: float = 0.7) -> Dict[str, Any]:
        """
        Generate text using Gemini API
        
        Args:
            prompt: Input text prompt
            max_length: Maximum length of generated text
            temperature: Sampling temperature
            
        Returns:
            Dict with generated text and metadata
        """
        try:
            config = types.GenerateContentConfig(
                temperature=temperature,
                max_output_tokens=max_length
            )
            
            response = self.client.models.generate_content(
                model=self.model_name,
                contents=prompt,
                config=config
            )
            
            return {
                "success": True,
                "response": response.text,
                "model": self.model_name,
                "device": "cloud"
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "model": self.model_name
            }
    
    def create_presentation_slides(self, topic: str, num_slides: int = 5) -> Dict[str, Any]:
        """
        Create presentation slides using Gemini AI with AI-generated images
        
        Args:
            topic: Topic for the presentation
            num_slides: Number of slides to create
            
        Returns:
            Dict with slide content and file path
        """
        try:
            # Generate slide content using Gemini
            prompt = f"""Tạo một bài thuyết trình về chủ đề: {topic}

Yêu cầu:
- Tạo {num_slides} slide
- Mỗi slide có tiêu đề và nội dung chi tiết
- Nội dung phải logic, mạch lạc
- Sử dụng bullet points
- Thêm trường "image_prompt" cho một số slide để tạo hình minh họa (bằng tiếng Anh, mô tả chi tiết)
- Format JSON như sau:
{{
  "title": "Tiêu đề bài thuyết trình",
  "slides": [
    {{
      "title": "Tiêu đề slide",
      "content": ["Điểm 1", "Điểm 2", "Điểm 3"],
      "image_prompt": "detailed image description in English (optional)"
    }}
  ]
}}

Chỉ trả về JSON, không thêm text khác."""

            config = types.GenerateContentConfig(
                temperature=0.7
            )
            
            response = self.client.models.generate_content(
                model=self.model_name,
                contents=prompt,
                config=config
            )
            
            # Parse JSON response
            text_response = response.text
            # Clean up markdown code blocks if present
            if "```json" in text_response:
                text_response = text_response.split("```json")[1].split("```")[0].strip()
            elif "```" in text_response:
                text_response = text_response.split("```")[1].split("```")[0].strip()
            
            slides_data = json.loads(text_response)
            
            # Initialize Text-to-Image model
            text_to_image = TextToImage()
            
            # Determine which slides should have images (about 60% of content slides)
            content_slides = slides_data.get("slides", [])
            num_images = min(int(num_slides * 0.6), len(content_slides))
            
            # Create PowerPoint presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = slides_data.get("title", topic)
            
            # Content slides
            images_generated = 0
            for idx, slide_data in enumerate(content_slides):
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                title_shape = slide.shapes.title
                title_shape.text = slide_data.get("title", "")
                
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                
                for point in slide_data.get("content", []):
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                
                # Generate and add image if needed
                should_add_image = (
                    images_generated < num_images and
                    (slide_data.get("image_prompt") or idx % 2 == 0)  # Add image to even slides or if prompt exists
                )
                
                if should_add_image:
                    # Use provided prompt or generate one from slide title
                    image_prompt = slide_data.get("image_prompt") or f"{slide_data.get('title', topic)}, professional illustration, high quality"
                    
                    # Generate image
                    result = text_to_image.generate_image(image_prompt)
                    
                    if result["success"]:
                        # Add image to slide (bottom right corner)
                        left = Inches(7)  # Right side
                        top = Inches(4.5)  # Bottom
                        width = Inches(2.5)
                        
                        slide.shapes.add_picture(
                            result["image_path"],
                            left, top,
                            width=width
                        )
                        images_generated += 1
                        print(f"✓ Added image to slide {idx + 1}")
            
            # Cleanup Text-to-Image model
            text_to_image.cleanup()
            
            # Save presentation
            os.makedirs("slides", exist_ok=True)
            filename = f"slides/presentation_{topic[:30].replace(' ', '_')}.pptx"
            prs.save(filename)
            
            return {
                "success": True,
                "filename": filename,
                "title": slides_data.get("title", topic),
                "num_slides": len(content_slides) + 1,
                "num_images": images_generated,
                "model": self.model_name
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "model": self.model_name
            }


# Global instance (will be initialized on first use)
local_llm = None
gemini_api = None

def get_local_llm() -> LocalLLM:
    """Get or create local LLM instance"""
    global local_llm
    if local_llm is None:
        local_llm = LocalLLM()
    return local_llm

def get_gemini_api(api_key: str, model_name: str = "gemini-2.5-flash") -> GeminiAPI:
    """Get or create Gemini API instance"""
    global gemini_api
    gemini_api = GeminiAPI(api_key, model_name)
    return gemini_api
