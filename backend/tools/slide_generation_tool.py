"""Automated Slide Generation Tool: Generate presentations from multiple documents using Gemini"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from google import genai
from google.genai import types
import os
from pathlib import Path
from typing import List, Dict, Any, Optional
import json
import base64
from PIL import Image
import io
import PyPDF2
import docx

class SlideGenerationTool:
    """Tool for generating presentation slides from documents"""
    
    def __init__(self, api_key: str):
        """Initialize the slide generation tool
        
        Args:
            api_key: Gemini API key
        """
        self.client = genai.Client(api_key=api_key)
        self.model_name = "gemini-3-flash-preview"
        
    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except Exception as e:
            print(f"Error extracting text from PDF: {e}")
            return ""
    
    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = docx.Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except Exception as e:
            print(f"Error extracting text from DOCX: {e}")
            return ""
    
    def extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            print(f"Error extracting text from TXT: {e}")
            return ""
    
    def extract_images_from_pdf(self, file_path: str, output_dir: str) -> List[str]:
        """Extract images from PDF file
        
        Args:
            file_path: Path to PDF file
            output_dir: Directory to save extracted images
            
        Returns:
            List of image file paths
        """
        try:
            images = []
            pdf_file = open(file_path, 'rb')
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    if '/XObject' not in page['/Resources']:
                        continue
                        
                    xObject = page['/Resources']['/XObject'].get_object()
                    
                    for obj in xObject:
                        try:
                            if xObject[obj]['/Subtype'] != '/Image':
                                continue
                            
                            width = xObject[obj]['/Width']
                            height = xObject[obj]['/Height']
                            
                            # Skip very small images (likely icons or artifacts)
                            if width < 50 or height < 50:
                                continue
                            
                            size = (width, height)
                            data = xObject[obj].get_data()
                            
                            # Skip if data is too small
                            if len(data) < 100:
                                continue
                            
                            # Save image
                            img_path = os.path.join(output_dir, f"extracted_page{page_num}_{obj[1:]}.png")
                            
                            # Try to open and convert the image
                            img = None
                            
                            # Method 1: Try direct color space conversion
                            if '/ColorSpace' in xObject[obj]:
                                color_space = xObject[obj]['/ColorSpace']
                                
                                if color_space == '/DeviceRGB':
                                    try:
                                        img = Image.frombytes('RGB', size, data)
                                    except:
                                        pass
                                elif color_space == '/DeviceCMYK':
                                    try:
                                        img = Image.frombytes('CMYK', size, data)
                                        img = img.convert('RGB')  # Convert CMYK to RGB
                                    except:
                                        pass
                                elif color_space == '/DeviceGray':
                                    try:
                                        img = Image.frombytes('L', size, data)
                                        img = img.convert('RGB')  # Convert grayscale to RGB
                                    except:
                                        pass
                            
                            # Method 2: Try to open from BytesIO
                            if img is None:
                                try:
                                    img = Image.open(io.BytesIO(data))
                                    # Convert to RGB if needed
                                    if img.mode not in ('RGB', 'RGBA'):
                                        img = img.convert('RGB')
                                except:
                                    pass
                            
                            # If we successfully got an image, save it
                            if img is not None:
                                # Ensure reasonable size
                                max_size = 2000
                                if img.width > max_size or img.height > max_size:
                                    img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
                                
                                # Save as PNG (RGB mode)
                                if img.mode == 'RGBA':
                                    # Create white background for transparent images
                                    background = Image.new('RGB', img.size, (255, 255, 255))
                                    background.paste(img, mask=img.split()[3] if len(img.split()) == 4 else None)
                                    img = background
                                elif img.mode != 'RGB':
                                    img = img.convert('RGB')
                                
                                img.save(img_path, 'PNG', optimize=True)
                                images.append(img_path)
                                
                        except Exception as e:
                            # Silently skip problematic images
                            continue
                            
                except Exception as e:
                    # Skip problematic pages
                    continue
            
            pdf_file.close()
            
            if images:
                print(f"✓ Successfully extracted {len(images)} image(s) from PDF")
            
            return images
            
        except Exception as e:
            print(f"Error extracting images from PDF: {e}")
            return []
    
    def extract_content_from_files(self, file_paths: List[str], images_dir: str) -> Dict[str, Any]:
        """Extract text and images from multiple files
        
        Args:
            file_paths: List of file paths to process
            images_dir: Directory to save extracted images
            
        Returns:
            Dictionary with combined text and image paths
        """
        combined_text = ""
        all_images = []
        
        for file_path in file_paths:
            file_ext = Path(file_path).suffix.lower()
            file_name = Path(file_path).name
            
            combined_text += f"\n\n--- Content from {file_name} ---\n\n"
            
            if file_ext == '.pdf':
                combined_text += self.extract_text_from_pdf(file_path)
                images = self.extract_images_from_pdf(file_path, images_dir)
                all_images.extend(images)
            elif file_ext == '.docx':
                combined_text += self.extract_text_from_docx(file_path)
            elif file_ext == '.txt':
                combined_text += self.extract_text_from_txt(file_path)
            elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
                # If it's an image file, add it to the images list
                all_images.append(file_path)
                combined_text += f"[Image file: {file_name}]\n"
            else:
                print(f"Unsupported file type: {file_ext}")
        
        return {
            "text": combined_text,
            "images": all_images
        }
    
    def analyze_with_gemini(self, content: Dict[str, Any], num_slides: int) -> Dict[str, Any]:
        """Analyze content and generate slide structure using Gemini
        
        Args:
            content: Dictionary with text and images
            num_slides: Target number of slides to generate
            
        Returns:
            Dictionary with slide structure
        """
        # Prepare the prompt for Gemini
        prompt = f"""You are an expert presentation designer. Analyze the following content from multiple documents and create a structured outline for a {num_slides}-slide presentation.

Content:
{content['text'][:15000]}  # Limit text to avoid token limits

Instructions:
1. Create exactly {num_slides} slides (including title slide)
2. Each slide should have:
   - A clear, concise title (max 8 words)
   - 3-5 bullet points summarizing key information
   - Optional: Suggest if an image would be relevant
3. Organize content logically with good flow
4. Make it visually appealing and professional
5. Focus on the most important information

Return your response as a JSON object with this structure:
{{
  "title": "Main presentation title",
  "slides": [
    {{
      "slide_number": 1,
      "title": "Slide title",
      "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
      "needs_image": true/false,
      "image_description": "Description of what image would fit"
    }}
  ]
}}

IMPORTANT: Return ONLY the JSON object, no additional text or markdown formatting."""

        try:
            response = self.client.models.generate_content(
                model=self.model_name,
                contents=prompt
            )
            
            # Extract JSON from response
            response_text = response.text.strip()
            
            # Remove markdown code blocks if present
            if response_text.startswith("```json"):
                response_text = response_text[7:]
            elif response_text.startswith("```"):
                response_text = response_text[3:]
            if response_text.endswith("```"):
                response_text = response_text[:-3]
            
            response_text = response_text.strip()
            
            # Parse JSON
            slide_structure = json.loads(response_text)
            return slide_structure
            
        except json.JSONDecodeError as e:
            print(f"Error parsing Gemini response: {e}")
            print(f"Response text: {response_text[:500]}")
            # Return a default structure
            return self._create_default_structure(content['text'], num_slides)
        except Exception as e:
            print(f"Error calling Gemini API: {e}")
            return self._create_default_structure(content['text'], num_slides)
    
    def _create_default_structure(self, text: str, num_slides: int) -> Dict[str, Any]:
        """Create a default slide structure if Gemini fails"""
        words = text.split()
        chunk_size = len(words) // (num_slides - 1) if num_slides > 1 else len(words)
        
        slides = [
            {
                "slide_number": 1,
                "title": "Document Summary",
                "content": ["Generated from uploaded documents"],
                "needs_image": False
            }
        ]
        
        for i in range(1, num_slides):
            start_idx = (i-1) * chunk_size
            end_idx = min(start_idx + chunk_size, len(words))
            chunk = " ".join(words[start_idx:end_idx])
            
            slides.append({
                "slide_number": i + 1,
                "title": f"Section {i}",
                "content": [chunk[:100] + "..."],
                "needs_image": False
            })
        
        return {
            "title": "Document Presentation",
            "slides": slides
        }
    
    def create_presentation(
        self,
        slide_structure: Dict[str, Any],
        images: List[str],
        output_path: str
    ) -> str:
        """Create PowerPoint presentation from slide structure
        
        Args:
            slide_structure: Slide structure from Gemini analysis
            images: List of extracted image paths
            output_path: Path to save the presentation
            
        Returns:
            Path to the created presentation
        """
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define modern color scheme
        title_color = RGBColor(31, 78, 121)  # Dark blue
        content_color = RGBColor(68, 68, 68)  # Dark gray
        accent_color = RGBColor(0, 176, 240)  # Bright blue
        
        # Create title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = slide_structure.get("title", "Document Summary")
        subtitle.text = "Auto-generated presentation"
        
        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        
        # Create content slides
        image_idx = 0
        for slide_info in slide_structure.get("slides", [])[1:]:  # Skip first as it's title
            # Use blank layout for more control
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5),
                Inches(9), Inches(1)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_info.get("title", "")
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = title_color
            
            # Add divider line
            line = slide.shapes.add_shape(
                1,  # Line shape
                Inches(0.5), Inches(1.4),
                Inches(9), Inches(0)
            )
            line.line.color.rgb = accent_color
            line.line.width = Pt(3)
            
            # Determine layout (with or without image)
            has_image = slide_info.get("needs_image", False) and image_idx < len(images)
            
            if has_image:
                # Two-column layout: text on left, image on right
                content_width = Inches(4.5)
                image_left = Inches(5.5)
                
                # Add image
                try:
                    img_path = images[image_idx]
                    slide.shapes.add_picture(
                        img_path,
                        image_left, Inches(2),
                        width=Inches(4), height=Inches(4.5)
                    )
                    image_idx += 1
                except Exception as e:
                    print(f"Error adding image: {e}")
                    content_width = Inches(9)  # Use full width if image fails
            else:
                content_width = Inches(9)
            
            # Add content
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2),
                content_width, Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for bullet_text in slide_info.get("content", []):
                p = content_frame.add_paragraph()
                p.text = bullet_text
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = content_color
                p.space_before = Pt(12)
                p.line_spacing = 1.2
        
        # Save presentation
        prs.save(output_path)
        return output_path
    
    def generate_slides_from_documents(
        self,
        file_paths: List[str],
        output_path: str,
        num_slides: int = 10,
        images_dir: Optional[str] = None
    ) -> Dict[str, Any]:
        """Main method to generate slides from documents
        
        Args:
            file_paths: List of document file paths
            output_path: Path to save the presentation
            num_slides: Target number of slides
            images_dir: Directory to save extracted images
            
        Returns:
            Dictionary with result information
        """
        try:
            # Set up images directory
            if images_dir is None:
                images_dir = os.path.join(os.path.dirname(output_path), "images")
            
            os.makedirs(images_dir, exist_ok=True)
            
            # Step 1: Extract content from files
            print("📄 Extracting content from documents...")
            content = self.extract_content_from_files(file_paths, images_dir)
            
            if not content['text'].strip():
                return {
                    "success": False,
                    "error": "No text content extracted from documents"
                }
            
            # Step 2: Analyze with Gemini
            print("🤖 Analyzing content with Gemini...")
            slide_structure = self.analyze_with_gemini(content, num_slides)
            
            # Step 3: Create presentation
            print("📊 Creating presentation...")
            result_path = self.create_presentation(
                slide_structure,
                content['images'],
                output_path
            )
            
            return {
                "success": True,
                "output_path": result_path,
                "num_slides": len(slide_structure.get("slides", [])),
                "num_images": len(content['images']),
                "title": slide_structure.get("title", "Document Summary")
            }
            
        except Exception as e:
            print(f"Error generating slides: {e}")
            return {
                "success": False,
                "error": str(e)
            }


def get_slide_generation_tool(api_key: str) -> SlideGenerationTool:
    """Factory function to create SlideGenerationTool instance
    
    Args:
        api_key: Gemini API key
        
    Returns:
        SlideGenerationTool instance
    """
    return SlideGenerationTool(api_key=api_key)
